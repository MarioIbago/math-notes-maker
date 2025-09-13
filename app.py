# app.py — Imagen/Text/PDF/PPTX ➜ Sheet Cheat en PDF (LaTeX)
# ===========================================================

import streamlit as st
from openai import OpenAI
from PIL import Image
from io import BytesIO
import base64
import subprocess
import os
import re
import tempfile
import PyPDF2
from pptx import Presentation
import shutil
import unicodedata

st.set_page_config(
    page_title="Sheet Cheat en PDF",
    page_icon="✏️",
    layout="centered",
    initial_sidebar_state="collapsed"  # 👈 fuerza sidebar colapsada
)


# ------------------ CONFIG ------------------
API_KEY = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY", "")
if not API_KEY:
    st.warning("⚠️ Configura tu OPENAI_API_KEY en Secrets (Streamlit) o como variable de entorno.")
client = OpenAI(api_key=API_KEY)

st.title("📘 Imagen / Texto / PDF / PPTX ➜ Sheet Cheat en PDF (by Mario Ibarra)")

# ------------------ PROMPT (más extenso) ------------------
PROMPT_CHEATSHEET = (
    "Actúas como un asistente experto en hojas de trucos (cheat sheets). "
    "Devuelve EXCLUSIVAMENTE un DOCUMENTO LaTeX completo y compilable. "
    "FORMATO OBLIGATORIO DE SALIDA: la PRIMERA línea debe ser exactamente 'COOR-BO-ZY'; "
    "a partir de la SEGUNDA línea comienza el documento LaTeX, SIN bloques de ``` ni texto fuera del LaTeX.\n\n"

    "PREÁMBULO (descríbelo, no pegues este texto): clase 'article' 11pt; español con babel; UTF-8 y T1; Latin Modern; "
    "amsmath, amssymb, amsthm, mathtools; geometry A4 (márgenes ~2–2.5 cm); enumitem; xcolor; tcolorbox sobrio; "
    "microtype; hyperref; graphicx. NO uses 'titlesec' ni \\titleformat/\\titlespacing; evita TikZ y paquetes no listados; "
    "nada de \\makeatletter, \\input, \\write18.\n\n"

    "INFIERE el <tema> a partir de la entrada (texto/imagen). Si el tema no está claro, elige el más probable y sé consistente.\n\n"

    "ESTRUCTURA EXACTA (con foco en FÓRMULAS, EXPLICACIONES y USO):\n"
    "1) Portada simple: \\title{Sheet Cheat: <tema>}, \\author{}, \\date{}, \\maketitle.\n"
    "2) \\section{Introducción}: 3–6 líneas (qué es, para qué sirve, contexto típico, supuestos básicos).\n"
    "3) \\section{Definición}: un tcolorbox con definición formal (incluye variables, dominios, unidades si aplica) "
    "   y al menos una ecuación en display.\n"
    "4) \\section{Propiedades}: 6–10 viñetas (1–2 líneas cada una) usando notación matemática; incluye condiciones "
    "   (dominio, continuidad, diferenciabilidad, linealidad, límites, casos particulares/pieza a pieza si procede).\n"
    "5) \\section{Fórmulas clave}: TODAS las fórmulas detectadas + las implícitas importantes que falten; agrúpalas por bloques "
    "   (comentarios de bloque) y usa equation* o align*; cuando aplique muestra equivalencias, formas alternativas y casos límite.\n"
    "6) \\section{Derivaciones mínimas}: 2–4 derivaciones limpias (≤6 líneas cada una) que muestren de dónde salen fórmulas clave; "
    "   evita pasos triviales; resalta el paso crítico o truco (por ejemplo, sustitución, factorización, regla de L'Hôpital, etc.).\n"
    "7) \\section{Ejemplos rápidos}: 2–3 ejemplos numéricos breves (2–6 líneas cada uno) resueltos paso a paso, "
    "   con verificación o interpretación final (unidades o significado).\n"
    "8) \\section{Aplicaciones}: enumerate de 4–7 ítems; cada ítem inicia con \\textbf{Etiqueta:} y una explicación breve y concreta.\n"
    "9) \\section{Símbolos y notación}: tabla/índice compacto de variables y constantes (1 línea por elemento).\n"
    "10) \\section{Errores comunes y buenas prácticas}: lista breve de 4–6 bullets (p. ej. signos, dominios, redondeos, unidades, orden de operaciones).\n"
    "11) \\section{Resumen de fórmulas esenciales}: tcolorbox con 5–10 fórmulas \"de oro\" en display; bajo cada una, "
    "   una nota de 1 línea (dominio/uso típico/alerta).\n\n"

    "ESTILO Y CALIDAD:\n"
    "• Español claro y conciso; objetivo 1–2 páginas. "
    "• Al final, añade una línea de pie de página en LaTeX que diga: Desarrollado por [MarioIbago](https://github.com/MarioIbago). "
    "• Matemáticas limpias: \\frac, potencias, sub/superscripts; alinea ecuaciones cuando ayude a la lectura. "
    "• Evita adornos innecesarios, imágenes, comentarios LaTeX o paquetes extra. Balancea todos los entornos; debe compilar en pdflatex.\n\n"

    "ENTRADA: recibirás texto o extracto de imagen; destila todas las fórmulas relevantes, explica brevemente cada bloque, "
    "añade notas de uso/errores comunes y cierra con un tcolorbox de fórmulas esenciales bien formateadas."
)

# ------------------ UTILS -------------------
_TITLESec_PAT = re.compile(r'\\usepackage\\s*\\{?\\s*titlesec\\s*\\}?', re.I)
_TITLEFORMAT_PAT = re.compile(r'\\title(?:format|spacing)\\*?[^\\n]*', re.I)

def _image_to_base64(uploaded_file):
    buf = BytesIO()
    img = Image.open(uploaded_file).convert("RGB")
    img.save(buf, format="JPEG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")

def extract_text_from_pdf(uploaded_file):
    try:
        reader = PyPDF2.PdfReader(uploaded_file)
        parts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                parts.append(page_text)
        return "\n".join(parts).strip()
    except Exception as e:
        st.error(f"Error leyendo PDF: {e}")
        return ""

def extract_text_from_pptx(uploaded_file):
    try:
        prs = Presentation(uploaded_file)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts).strip()
    except Exception as e:
        st.error(f"Error leyendo PPTX: {e}")
        return ""

def sanitize_latex(content: str) -> str:
    lines = []
    for line in content.splitlines():
        s = line.strip()
        if s.startswith("```") or s.startswith("% !TEX"):
            continue
        lines.append(line)
    txt = "\n".join(lines)
    txt = txt.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\r\n", "\n")
    txt = txt.replace("\ufeff", "").strip()
    if txt.startswith("COOR-BO-ZY"):
        txt = txt.split("\n", 1)[-1] if "\n" in txt else ""
    txt = _TITLESec_PAT.sub("", txt)
    txt = _TITLEFORMAT_PAT.sub("", txt)
    return txt.strip()

def ensure_full_document(latex_code: str) -> str:
    if "\\begin{document}" in latex_code:
        return latex_code
    wrapper = (
        "\\documentclass[11pt]{article}\n"
        "\\usepackage[spanish, es-tabla]{babel}\n"
        "\\usepackage[utf8]{inputenc}\n"
        "\\usepackage[T1]{fontenc}\n"
        "\\usepackage{lmodern}\n"
        "\\usepackage{amsmath,amssymb,amsthm,mathtools}\n"
        "\\usepackage[a4paper,margin=2.0cm]{geometry}\n"
        "\\usepackage{enumitem}\n"
        "\\usepackage{xcolor}\n"
        "\\usepackage{tcolorbox}\n"
        "\\usepackage{graphicx}\n"
        "\\usepackage{hyperref}\n"
        "\\usepackage{microtype}\n"
        "\\tcbset{colback=gray!3,colframe=black!50,boxrule=0.5pt,arc=2pt}\n\n"
        "\\begin{document}\n" + latex_code + "\n\\end{document}\n"
    )
    return wrapper

def call_openai(prompt, notes_text=None, image_b64=None):
    if not API_KEY:
        st.error("No hay API key configurada.")
        return ""
    try:
        if image_b64:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                max_tokens=3000,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt + "\n\nTexto de entrada: (extraído de imagen)"},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_b64}"}}
                    ]
                }]
            )
        else:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                max_tokens=3000,
                messages=[{"role": "user", "content": prompt + f"\n\nTexto de entrada:\n{notes_text or ''}"}]
            )
        content = (resp.choices[0].message.content or "").strip()
    except Exception as e:
        st.error(f"Error al llamar a OpenAI: {e}")
        return ""

    content = sanitize_latex(content)
    content = ensure_full_document(content)
    return content

# --------- Detección de tema y nombre de archivo ----------
TITLE_RE = re.compile(r"\\title\\{\\s*Sheet\\s*Cheat:\\s*([^}]*)\\}", re.IGNORECASE)

def _slugify(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9]+", "-", text).strip("-").lower()
    return text or "tema"

def infer_topic(latex_code: str, notes_text: str, sidebar_topic: str) -> str:
    # prioridad: topic manual en sidebar > \title{Sheet Cheat: <tema>} > primera línea entrada > 'tema'
    if sidebar_topic and sidebar_topic.strip():
        return sidebar_topic.strip()
    m = TITLE_RE.search(latex_code or "")
    if m:
        t = m.group(1).strip()
        if t:
            return t
    if notes_text:
        first = notes_text.strip().splitlines()[0]
        if first and len(first) <= 80:
            return first.strip()
    return "Tema"

def build_filenames(topic: str):
    slug = _slugify(topic)
    base = f"sheat_cheat_{slug}"
    return base + ".tex", base + ".pdf"

# ------------------ Compilación (con fallback opcional) ------------------
def _has_cmd(cmd: str) -> bool:
    return shutil.which(cmd) is not None

def compile_pdf(latex_code: str, engine_preference=("pdflatex", "tectonic")):
    if not latex_code.strip():
        st.error("No hay código LaTeX para compilar.")
        return None

    with tempfile.TemporaryDirectory() as tmpdir:
        tex_path = os.path.join(tmpdir, "sheat_cheat.tex")
        with open(tex_path, "w", encoding="utf-8") as f:
            f.write(latex_code)

        for engine in engine_preference:
            if engine == "pdflatex" and _has_cmd("pdflatex"):
                for _ in range(2):
                    result = subprocess.run(
                        ["pdflatex", "-interaction=nonstopmode", "-halt-on-error", "-file-line-error", "sheat_cheat.tex"],
                        cwd=tmpdir,
                        stdout=subprocess.PIPE, stderr=subprocess.PIPE
                    )
                    if result.returncode != 0:
                        log_path = os.path.join(tmpdir, "sheat_cheat.log")
                        stdout = result.stdout.decode(errors="ignore")
                        stderr = result.stderr.decode(errors="ignore")
                        logtxt = ""
                        if os.path.exists(log_path):
                            with open(log_path, "r", encoding="utf-8", errors="ignore") as logf:
                                logtxt = logf.read()
                        st.text_area("❌ Error en LaTeX (pdflatex)", logtxt or stdout or stderr, height=360)
                        return None
                pdf_path = os.path.join(tmpdir, "sheat_cheat.pdf")
                if os.path.exists(pdf_path):
                    with open(pdf_path, "rb") as fpdf:
                        return fpdf.read()
                st.error("❌ No se encontró el PDF tras compilar con pdflatex.")
                return None

            if engine == "tectonic" and _has_cmd("tectonic"):
                result = subprocess.run(
                    ["tectonic", "sheat_cheat.tex", "--keep-logs"],
                    cwd=tmpdir,
                    stdout=subprocess.PIPE, stderr=subprocess.PIPE
                )
                if result.returncode != 0:
                    stdout = result.stdout.decode(errors="ignore")
                    stderr = result.stderr.decode(errors="ignore")
                    log_path = os.path.join(tmpdir, "sheat_cheat.log")
                    logtxt = ""
                    if os.path.exists(log_path):
                        with open(log_path, "r", encoding="utf-8", errors="ignore") as logf:
                            logtxt = logf.read()
                    st.text_area("❌ Error en LaTeX (tectonic)", logtxt or stdout or stderr, height=360)
                    return None
                pdf_path = os.path.join(tmpdir, "sheat_cheat.pdf")
                if os.path.exists(pdf_path):
                    with open(pdf_path, "rb") as fpdf:
                        return fpdf.read()
                st.error("❌ No se encontró el PDF tras compilar con Tectonic.")
                return None

        st.error(
            "No se encontró un motor LaTeX. Instala TeX Live o Tectonic con `packages.txt`:\n"
            "- texlive-latex-base\n- texlive-latex-recommended\n- texlive-latex-extra\n"
            "- texlive-fonts-recommended\n- texlive-lang-spanish\n- (opcional) tectonic"
        )
        return None

# ------------------ INTERFAZ ------------------
if False:
    with st.sidebar:
        st.header("⚙️ Opciones")
        sidebar_topic = st.text_input("Tema (opcional, sobrescribe el detectado)", value="")
        st.markdown(
            "- El nombre del archivo usará: `sheat_cheat_<tema>`\n"
            "- Si no escribes tema, se detecta del LaTeX o de tu entrada."
        )
        st.divider()
        st.caption("Usa `packages.txt` en Streamlit Cloud para instalar LaTeX. Opcional: agrega `tectonic`.")
        
        if False:  # 👈 nunca se ejecuta
        with st.sidebar:
            st.header("⚙️ Opciones")
            sidebar_topic = st.text_input("Tema (opcional, sobrescribe el detectado)", value="")

if mode == "Subir imagen":
    up = st.file_uploader("📤 Sube una imagen (JPG/PNG)", type=["jpg","jpeg","png"])
    if up is not None:
        st.image(up, caption="Imagen cargada", width=300)
        image_b64 = _image_to_base64(up)

elif mode == "Escribir texto":
    notes_text = st.text_area("✍️ Escribe o pega tus notas", height=220, placeholder="Tema o contenido...")

elif mode == "Subir PDF":
    up = st.file_uploader("📤 Sube un PDF", type=["pdf"])
    if up is not None:
        notes_text = extract_text_from_pdf(up)
        st.text_area("📄 Texto extraído del PDF", notes_text, height=220)

elif mode == "Subir PPTX":
    up = st.file_uploader("📤 Sube un PPTX", type=["pptx"])
    if up is not None:
        notes_text = extract_text_from_pptx(up)
        st.text_area("📊 Texto extraído del PPTX", notes_text, height=220)

st.divider()

if st.button("⚡ Generar Sheet Cheat", use_container_width=True):
    st.info("⏳ Generando LaTeX con GPT...")
    latex_code = call_openai(PROMPT_CHEATSHEET, notes_text=notes_text, image_b64=image_b64)

    if not latex_code:
        st.error("No se obtuvo contenido de LaTeX.")
        st.stop()

    # Detectar tema y construir nombres de archivo
    topic = infer_topic(latex_code, notes_text, sidebar_topic)
    tex_name, pdf_name = build_filenames(topic)

    st.success(f"🧾 Tema detectado: **{topic}**")
    st.caption(f"Nombres de archivo: `{tex_name}` y `{pdf_name}`")

    st.subheader("📄 Código LaTeX completo")
    st.code(latex_code, language="latex")

    st.download_button(
        f"📥 Descargar {tex_name}",
        data=latex_code.encode("utf-8"),
        file_name=tex_name,
        mime="text/plain",
        use_container_width=True
    )

    st.info("🛠️ Compilando (pdflatex ➜ fallback a tectonic si está disponible)...")
    pdf_bytes = compile_pdf(latex_code)

    if pdf_bytes:
        st.success("✅ PDF generado correctamente.")
        st.download_button(
            f"📥 Descargar {pdf_name}",
            data=pdf_bytes,
            file_name=pdf_name,
            mime="application/pdf",
            use_container_width=True
        )
    else:
        st.error("❌ Falló la compilación del PDF. Revisa el LaTeX o el log mostrado.")
